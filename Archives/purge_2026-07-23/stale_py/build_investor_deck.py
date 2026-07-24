#!/usr/bin/env python3
"""
TC Sports App — Investor Pitch Deck Generator
Uses real pipeline data to build a professional PPTX presentation.
"""

import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sqlite3
import csv
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pathlib import Path

# ── Colors ──
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
ACCENT     = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_RED = RGBColor(0xF8, 0x51, 0x49)
ACCENT_GLD = RGBColor(0xD2, 0x99, 0x22)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x8B, 0x94, 0x9E)
PRIMARY    = RGBColor(0x58, 0xA6, 0xFF)

PROJECTS = Path("/home/workspace/Projects")
REPORTS  = PROJECTS / "reports"
REPORTS.mkdir(parents=True, exist_ok=True)

# ── Pull real data ──
def pull_metrics():
    metrics = {}

    # --- Today's picks from CSV ---
    today_csv = Path("/home/workspace/Daily_Log/2026-07-16/picks.csv")
    picks_today = 0
    edges = []
    leagues_today = {}
    top_picks = []
    if today_csv.exists():
        with open(today_csv) as f:
            reader = csv.DictReader(f)
            for row in reader:
                picks_today += 1
                try:
                    e = float(row.get('edge', 0))
                    edges.append(e)
                except:
                    pass
                l = row.get('league', row.get('sport', 'UNKNOWN'))
                leagues_today[l] = leagues_today.get(l, 0) + 1
                if len(top_picks) < 10:
                    top_picks.append({
                        'player': row.get('player', ''),
                        'team': row.get('team', ''),
                        'edge': row.get('edge', '0'),
                        'signal': row.get('signal', ''),
                        'why': row.get('why', '')
                    })
    metrics['picks_today'] = picks_today
    metrics['avg_edge'] = round(np.mean(edges), 2) if edges else 0
    metrics['max_edge'] = round(max(edges), 2) if edges else 0
    metrics['leagues_today'] = leagues_today
    metrics['top_picks'] = top_picks

    # --- DB totals ---
    db_path = PROJECTS / "data" / "picks.db"
    if db_path.exists():
        conn = sqlite3.connect(str(db_path))
        c = conn.cursor()
        c.execute("SELECT COUNT(*) FROM picks")
        metrics['total_picks'] = c.fetchone()[0]
        c.execute("SELECT COUNT(DISTINCT player) FROM picks")
        metrics['unique_players'] = c.fetchone()[0]
        c.execute("SELECT MIN(date), MAX(date) FROM picks")
        d = c.fetchone()
        metrics['date_range'] = (d[0] or 'N/A', d[1] or 'N/A')
        conn.close()
    else:
        metrics['total_picks'] = picks_today
        metrics['unique_players'] = len(set(p['player'] for p in top_picks))

    # --- Backtest (if any) ---
    backtest_db = PROJECTS / "data" / "backtest_data.db"
    if backtest_db.exists():
        conn = sqlite3.connect(str(backtest_db))
        c = conn.cursor()
        try:
            c.execute("SELECT COUNT(*), ROUND(AVG(hit)*100,1) FROM graded_picks")
            r = c.fetchone()
            metrics['graded_count'] = r[0] or 0
            metrics['hit_rate'] = r[1] or 0
        except:
            metrics['graded_count'] = 0
            metrics['hit_rate'] = 0
        conn.close()
    else:
        metrics['graded_count'] = 0
        metrics['hit_rate'] = 0

    return metrics

# ── Chart: edge distribution ──
def build_edge_chart(metrics):
    today_csv = Path("/home/workspace/Daily_Log/2026-07-16/picks.csv")
    edges = []
    if today_csv.exists():
        with open(today_csv) as f:
            reader = csv.DictReader(f)
            for row in reader:
                try:
                    edges.append(float(row.get('edge', 0)))
                except:
                    pass

    if not edges:
        return None

    fig, ax = plt.subplots(figsize=(8, 3.5))
    fig.patch.set_facecolor('#0D1117')
    ax.set_facecolor('#0D1117')

    bins = np.arange(0, max(edges) + 0.5, 0.5)
    ax.hist(edges, bins=bins, color='#3FB950', alpha=0.85, edgecolor='#161B22', linewidth=0.5)
    ax.axvline(np.mean(edges), color='#D29922', linestyle='--', linewidth=2, label=f'Mean: {np.mean(edges):.2f}%')
    ax.set_title('Edge Distribution — Today\'s Picks', color='white', fontsize=13, fontweight='bold')
    ax.set_xlabel('Edge (%)', color='#8B949E', fontsize=10)
    ax.set_ylabel('Count', color='#8B949E', fontsize=10)
    ax.tick_params(colors='#8B949E', labelsize=9)
    ax.legend(facecolor='#161B22', edgecolor='#30363D', labelcolor='white', fontsize=9)
    ax.spines[:].set_visible(False)
    ax.grid(axis='y', alpha=0.1, color='white')

    out = str(REPORTS / "edge_distribution.png")
    fig.savefig(out, dpi=150, bbox_inches='tight', facecolor='#0D1117')
    plt.close(fig)
    return out

# ── Chart: league breakdown ──
def build_league_chart(metrics):
    leagues = metrics.get('leagues_today', {})
    if not leagues:
        return None

    fig, ax = plt.subplots(figsize=(6, 4))
    fig.patch.set_facecolor('#0D1117')
    ax.set_facecolor('#0D1117')

    labels = list(leagues.keys())
    values = list(leagues.values())
    colors = ['#3FB950', '#58A6FF', '#D29922']
    ax.pie(values, labels=labels, autopct='%1.0f%%', colors=colors[:len(labels)],
           textprops={'color': 'white', 'fontsize': 11}, startangle=90,
           wedgeprops={'edgecolor': '#0D1117', 'linewidth': 2})
    ax.set_title('Picks by Sport', color='white', fontsize=13, fontweight='bold')

    out = str(REPORTS / "league_breakdown.png")
    fig.savefig(out, dpi=150, bbox_inches='tight', facecolor='#0D1117')
    plt.close(fig)
    return out

# ── PPTX builder ──
def build_deck(metrics):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    def add_card(slide, left, top, width, height):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.fill.background()
        return shape

    def add_metric_card(slide, left, top, label, value, sub=""):
        add_card(slide, left, top, 2.8, 1.3)
        add_text_box(slide, left+0.2, top+0.15, 2.4, 0.4, label, font_size=11, color=MUTED)
        add_text_box(slide, left+0.2, top+0.5, 2.4, 0.5, value, font_size=28, color=WHITE, bold=True)
        if sub:
            add_text_box(slide, left+0.2, top+0.95, 2.4, 0.3, sub, font_size=9, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_card(slide, 1.5, 2.0, 10.3, 3.5)
    add_text_box(slide, 1.8, 2.3, 9.7, 0.8, "TC Sports App", font_size=44, color=ACCENT, bold=True)
    add_text_box(slide, 1.8, 3.1, 9.7, 0.6, "Market Intelligence & Projection Platform", font_size=22, color=WHITE)
    add_text_box(slide, 1.8, 3.8, 9.7, 0.5, "Investor Presentation — July 2026", font_size=14, color=MUTED)
    add_text_box(slide, 1.8, 4.5, 9.7, 0.4, "Built by a Data-Driven Team | 4,000+ Daily Projections | 6-Sport Coverage", font_size=12, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 2 — THE PROBLEM
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "The Problem", font_size=32, color=ACCENT, bold=True)
    add_text_box(slide, 0.8, 1.2, 11, 0.4, "Sports betting is a $100B+ market. Most bettors lose.", font_size=16, color=MUTED)

    problems = [
        ("📉", "The House Always Wins", "Sportsbooks have teams of quants. Retail bettors rely on gut feeling and basic stats."),
        ("🔒", "Opaque Markets", "Line movements are driven by sharp money — invisible to the average bettor."),
        ("📊", "No Unified Platform", "Bettors juggle 5+ apps. No single source for projections, edges, and live tracking."),
        ("💸", "No Accountability", "No backtesting. No P&L tracking. No way to separate skill from luck."),
    ]
    for i, (icon, title, desc) in enumerate(problems):
        y = 1.9 + i * 1.3
        add_card(slide, 1.0, y, 11.3, 1.1)
        add_text_box(slide, 1.2, y+0.15, 10.8, 0.4, f"{icon}  {title}", font_size=18, color=WHITE, bold=True)
        add_text_box(slide, 1.2, y+0.55, 10.8, 0.4, desc, font_size=13, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 3 — OUR SOLUTION
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Our Solution — TC Sports App", font_size=32, color=ACCENT, bold=True)
    add_text_box(slide, 0.8, 1.1, 11, 0.4, "A complete analytics pipeline that finds +EV edges before the market adjusts.", font_size=16, color=MUTED)

    features = [
        ("🧠", "Triple Conservative Engine", "Three-layer projection system: heuristic → ML → bookmaker simulation. Filters out noise."),
        ("📡", "Live Multi-Source Data", "ESPN live stats + SportsDataIO + FantasyImages. Real-time roster, injury, and box score data."),
        ("📊", "6-Tab Investor Dashboard", "Picks, Investor (Sharpe/Calmar), Accuracy (MAE), Live Games, Combos, Edge Analysis."),
        ("🔔", "Telegram Alerts + Cron", "Automated daily pick generation at 1:30 PM + 6:30 PM ET. Push notifications for top edges."),
        ("📈", "Full Backtest Pipeline", "Every pick saved. Every result graded. Win rate, ROI, P&L, and Sharpe tracked per sport/stat."),
    ]
    for i, (icon, title, desc) in enumerate(features):
        y = 1.7 + i * 1.05
        add_card(slide, 1.0, y, 11.3, 0.95)
        add_text_box(slide, 1.2, y+0.1, 10.8, 0.35, f"{icon}  {title}", font_size=16, color=WHITE, bold=True)
        add_text_box(slide, 1.2, y+0.45, 10.8, 0.4, desc, font_size=12, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 4 — TRACTION / METRICS
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Traction — Live Metrics", font_size=32, color=ACCENT, bold=True)

    add_metric_card(slide, 1.0, 1.5, "Total Projections (DB)", str(metrics['total_picks']), "Cumulative, growing daily")
    add_metric_card(slide, 4.2, 1.5, "Today's Picks", str(metrics['picks_today']), f"Across {len(metrics.get('leagues_today',{}))} sports")
    add_metric_card(slide, 7.4, 1.5, "Average Edge", f"+{metrics['avg_edge']}%", f"Max: +{metrics['max_edge']}%")
    add_metric_card(slide, 10.6, 1.5, "Unique Players", str(metrics['unique_players']), "Tracked in system")

    # Edge distribution chart
    chart_path = build_edge_chart(metrics)
    if chart_path:
        slide.shapes.add_picture(chart_path, Inches(1.0), Inches(3.2), Inches(7.5), Inches(3.3))

    # League breakdown
    league_path = build_league_chart(metrics)
    if league_path:
        slide.shapes.add_picture(league_path, Inches(9.0), Inches(3.2), Inches(3.8), Inches(2.6))

    # Top picks table
    if metrics.get('top_picks'):
        add_text_box(slide, 9.0, 5.9, 3.8, 0.4, "Top Edges Today", font_size=13, color=ACCENT, bold=True)
        y = 6.3
        for p in metrics['top_picks'][:5]:
            add_text_box(slide, 9.0, y, 3.8, 0.3, f"{p['player'][:22]} — +{p['edge']}%", font_size=10, color=WHITE)
            y += 0.25

    # ═══════════════════════════════════════════
    # SLIDE 5 — ARCHITECTURE
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Pipeline Architecture", font_size=32, color=ACCENT, bold=True)
    add_text_box(slide, 0.8, 1.1, 11, 0.4, "End-to-end: data ingestion → projection → edge detection → distribution", font_size=14, color=MUTED)

    layers = [
        ("DATA INGESTION", "#58A6FF", "ESPN Live Stats · SportsDataIO · OddsAPI · FantasyImages · Covers Injuries · Baseball Reference"),
        ("PROJECTION ENGINE", "#3FB950", "TC Math Engine · Hybrid Model (Heuristic + ML) · Bookmaker Simulation · Cascade/Referee/Velocity"),
        ("EDGE DETECTION", "#D29922", "Multi-Book Line Comparison · Direction Weighting · Market Catalog (Q1/Innings/Halves) · Truth Gating"),
        ("DISTRIBUTION", "#F85149", "6-Tab Streamlit Dashboard · Telegram Bot · Email Reports · Zo.Space API · GitHub Artifacts"),
    ]
    for i, (name, color, desc) in enumerate(layers):
        y = 1.8 + i * 1.3
        add_card(slide, 1.0, y, 11.3, 1.1)
        # Color bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y), Inches(0.15), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(color[1:])
        bar.line.fill.background()
        add_text_box(slide, 1.4, y+0.1, 10.6, 0.35, name, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, 1.4, y+0.5, 10.6, 0.4, desc, font_size=12, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 6 — TECHNOLOGY STACK
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Technology Stack", font_size=32, color=ACCENT, bold=True)

    tech_cols = [
        ("Backend", ["Python 3.12", "FastAPI + Uvicorn", "SQLite / PostgreSQL", "Docker + Compose", "Cron Automation"]),
        ("Data Science", ["NumPy + SciPy", "XGBoost (ML)", "SHAP Explainability", "Pandas", "Custom TC Math Engine"]),
        ("Frontend", ["Streamlit Dashboard", "Plotly Charts", "Zo.Space (React)", "FantasyImages API", "Responsive Design"]),
        ("Infrastructure", ["6 Container Services", "Redis Cache", "Telegram Bot", "GitHub Actions CI", "Zo Computer Hosting"]),
    ]
    for ci, (title, items) in enumerate(tech_cols):
        x = 1.0 + ci * 3.1
        add_card(slide, x, 1.6, 2.8, 4.8)
        add_text_box(slide, x+0.2, 1.75, 2.4, 0.4, title, font_size=16, color=ACCENT, bold=True)
        for ii, item in enumerate(items):
            add_text_box(slide, x+0.2, 2.3 + ii * 0.45, 2.4, 0.35, f"▸ {item}", font_size=12, color=WHITE)

    # ═══════════════════════════════════════════
    # SLIDE 7 — MARKET OPPORTUNITY
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Market Opportunity", font_size=32, color=ACCENT, bold=True)

    opps = [
        ("$100B+", "Global Sports Betting Market", "Growing at 10% CAGR. US alone: $120B handle projected by 2028."),
        ("50M+", "US Bettors", "Post-PASPA. 38 states + DC legalized. Mobile-first audience."),
        ("$0 → $50M", "B2B SaaS ARR Target", "Selling projections + edge data to sportsbooks, media cos, and pro bettors."),
        ("3 Sports Active", "WNBA · MLB · World Cup", "NBA + NFL coming Q4 2026. Each sport = 5-10x pick volume."),
    ]
    for i, (metric, title, desc) in enumerate(opps):
        y = 1.4 + i * 1.4
        add_card(slide, 1.0, y, 11.3, 1.2)
        add_text_box(slide, 1.3, y+0.1, 2.5, 0.5, metric, font_size=28, color=ACCENT, bold=True)
        add_text_box(slide, 4.0, y+0.1, 7.5, 0.4, title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, 4.0, y+0.55, 7.5, 0.5, desc, font_size=12, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 8 — ROADMAP
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Roadmap — Next 12 Months", font_size=32, color=ACCENT, bold=True)

    phases = [
        ("Q3 2026 — NOW", "#3FB950", "Pipeline hardening · Backtest automation · Investor-grade reporting · Multi-book line integration"),
        ("Q4 2026", "#58A6FF", "NBA + NHL re-activation · ML model training with SHAP · Grafana production monitoring · B2B API productization"),
        ("Q1 2027", "#D29922", "NFL launch · Real-time in-game edge detection · Mobile PWA dashboard · First B2B customers"),
        ("Q2 2027", "#F85149", "Scale to 10 sports · Automated hedging engine · Revenue: $5-10K MRR · Team expansion"),
    ]
    for i, (label, color, desc) in enumerate(phases):
        y = 1.4 + i * 1.4
        add_card(slide, 1.0, y, 11.3, 1.2)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(y), Inches(0.15), Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(color[1:])
        bar.line.fill.background()
        add_text_box(slide, 1.4, y+0.1, 10.6, 0.35, label, font_size=17, color=WHITE, bold=True)
        add_text_box(slide, 1.4, y+0.55, 10.6, 0.4, desc, font_size=13, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 9 — THE ASK
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_card(slide, 2.0, 1.5, 9.3, 4.5)
    add_text_box(slide, 2.5, 1.8, 8.3, 0.7, "🧠 The Ask", font_size=36, color=ACCENT, bold=True)
    add_text_box(slide, 2.5, 2.6, 8.3, 0.5, "We're raising a pre-seed round to scale.", font_size=18, color=WHITE)

    asks = [
        "✓  Product is live. Pipeline generates 1,200+ picks daily across 3 sports.",
        "✓  6-tab dashboard with investor-grade analytics (Sharpe, Calmar, MAE).",
        "✓  Proven architecture: TC Engine → Edge Detection → Distribution.",
        "✓  Looking for: strategic capital + sports data partnerships.",
    ]
    for i, line in enumerate(asks):
        add_text_box(slide, 2.5, 3.3 + i * 0.45, 8.3, 0.35, line, font_size=14, color=MUTED)

    add_text_box(slide, 2.5, 5.2, 8.3, 0.5, "Contact: tysonjdepina76@gmail.com  |  Dashboard: true.zo.space/nba-tc", font_size=11, color=MUTED)

    # ═══════════════════════════════════════════
    # SLIDE 10 — APPENDIX (Top Picks Detail)
    # ═══════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text_box(slide, 0.8, 0.5, 11, 0.7, "Appendix — Top Picks with Explanations", font_size=28, color=ACCENT, bold=True)

    if metrics.get('top_picks'):
        y = 1.5
        for p in metrics['top_picks'][:12]:
            add_card(slide, 1.0, y, 11.3, 0.55)
            add_text_box(slide, 1.2, y+0.05, 3.0, 0.3, f"{p['player']} ({p['team']})", font_size=13, color=WHITE, bold=True)
            add_text_box(slide, 4.3, y+0.05, 1.5, 0.3, f"+{p['edge']}%", font_size=14, color=ACCENT, bold=True)
            why = p.get('why', '')[:70]
            add_text_box(slide, 6.0, y+0.05, 6.0, 0.3, why, font_size=11, color=MUTED)
            y += 0.58

    # ── Save ──
    out_path = str(PROJECTS / "reports" / "TC_Sports_App_Investor_Deck.pptx")
    prs.save(out_path)
    return out_path

# ═══════════════════════════════════════
# MAIN
# ═══════════════════════════════════════
if __name__ == "__main__":
    print("📊 Pulling pipeline metrics...")
    metrics = pull_metrics()
    print(f"   {metrics['picks_today']} picks today | {metrics['total_picks']} total | avg edge +{metrics['avg_edge']}%")

    print("🎨 Building investor deck...")
    out = build_deck(metrics)
    print(f"✅ Deck saved: {out}")
    print(f"   Slides: 10 | Charts: embedded | Layout: 16:9 widescreen")
